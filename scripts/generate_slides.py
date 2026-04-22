from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "slides"
OUT_PATH = OUT_DIR / "rr-quantum-simulator-deck.pptx"

BG = RGBColor(248, 249, 252)
NAVY = RGBColor(25, 45, 85)
BLUE = RGBColor(61, 106, 176)
TEAL = RGBColor(39, 119, 117)
GOLD = RGBColor(201, 155, 59)
TEXT = RGBColor(42, 48, 60)
MUTED = RGBColor(99, 109, 125)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(231, 236, 244)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.17), Inches(8.5), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(8.95), Inches(0.20), Inches(3.7), Inches(0.35))
        p = sub_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(214, 223, 238)


def add_bullets(slide, left, top, width, height, bullets, font_size=22, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(11)
    return box


def add_footer(slide, text="Principles of Operating Systems Project"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.08), Inches(6.0), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(10.5), Inches(1.8))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Interactive Round Robin\nQuantum Tuning Simulator"
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NAVY

    sub = slide.shapes.add_textbox(Inches(0.95), Inches(3.45), Inches(9.8), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = (
        "A terminal-based experiment platform for studying how time quantum selection "
        "changes responsiveness, overhead, and overall scheduling performance."
    )
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT

    info = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.15), Inches(4.0), Inches(1.15)
    )
    info.fill.solid()
    info.fill.fore_color.rgb = WHITE
    info.line.color.rgb = LIGHT
    tf = info.text_frame
    tf.paragraphs[0].text = "Mykhailo Chuvik"
    tf.paragraphs[0].font.name = "Aptos"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    p = tf.add_paragraph()
    p.text = "Principles of Operating Systems"
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED

    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(5.08), Inches(3.9), Inches(0.75)
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BLUE
    ribbon.line.fill.background()
    tf = ribbon.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = "Scheduling Policy Analysis"
    tf.paragraphs[0].font.name = "Aptos"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    add_footer(slide)


def section_slide(prs, title, bullets, subtitle=None, right_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, title, subtitle)
    add_bullets(slide, Inches(0.75), Inches(1.35), Inches(7.1), Inches(5.4), bullets)

    if right_note:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.5), Inches(4.2), Inches(4.8)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LIGHT
        tf = panel.text_frame
        tf.clear()
        for idx, line in enumerate(right_note):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Aptos"
            p.font.size = Pt(18 if idx == 0 else 15)
            p.font.bold = idx == 0
            p.font.color.rgb = NAVY if idx == 0 else TEXT
            p.space_after = Pt(8)

    add_footer(slide)


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Implementation Overview", "Core simulator structure")

    items = [
        ("Workload Parser", "Reads PID, arrival time, and CPU burst input files"),
        ("Ready Queue", "Maintains RR dispatch order as processes arrive and requeue"),
        ("Simulation Engine", "Advances time, dispatches jobs, charges context-switch cost"),
        ("Metrics Collector", "Computes wait, turnaround, response, throughput, utilization"),
        ("Visualization Layer", "Trace, HUD animation, and theatre animation modes"),
        ("Experiment Modes", "Single-run analysis, sweep mode, FCFS baseline, workload generator"),
    ]

    y = 1.45
    for idx, (head, body) in enumerate(items):
        x = 0.75 if idx < 3 else 6.9
        local_y = y + (idx % 3) * 1.65
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(local_y), Inches(5.45), Inches(1.15)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT
        tf = card.text_frame
        tf.paragraphs[0].text = head
        tf.paragraphs[0].font.name = "Aptos"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY
        p = tf.add_paragraph()
        p.text = body
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT

    add_footer(slide)


def metrics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Metrics and Experiment Design", "What the simulator measures")

    left = [
        "Average waiting time: how long processes spend waiting in READY",
        "Average turnaround time: arrival to completion",
        "Average response time: arrival to first CPU dispatch",
        "Context switches and total completion time",
        "Throughput and CPU utilization",
    ]
    add_bullets(slide, Inches(0.8), Inches(1.4), Inches(6.1), Inches(4.4), left, font_size=20)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.55), Inches(5.0), Inches(4.6)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LIGHT
    tf = panel.text_frame
    tf.clear()
    lines = [
        "Experimental variables",
        "Quantum range: sweep across multiple q values",
        "Context-switch cost: compare cs = 0, 1, 2",
        "Workload type: demo, interactive, mixed, generated",
        "Baseline comparison: Round Robin vs FCFS",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(18 if idx == 0 else 15)
        p.font.bold = idx == 0
        p.font.color.rgb = NAVY if idx == 0 else TEXT
        p.space_after = Pt(8)

    add_footer(slide)


def workloads_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Workloads Used", "Representative test scenarios")

    cards = [
        ("demo.txt", "Small four-process example for readable traces and quick sanity checks"),
        ("interactive-bursty.txt", "Many short jobs arriving frequently to emphasize response time"),
        ("big-mix.txt", "Larger mixed workload with short, medium, and long CPU bursts"),
        ("generated workloads", "Repeatable synthetic inputs using a seed for experiments"),
    ]
    x_positions = [0.75, 6.8, 0.75, 6.8]
    y_positions = [1.55, 1.55, 3.85, 3.85]
    colors = [BLUE, TEAL, GOLD, NAVY]
    for idx, (title, body) in enumerate(cards):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x_positions[idx]),
            Inches(y_positions[idx]),
            Inches(5.6),
            Inches(1.75),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT

        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x_positions[idx] + 0.2),
            Inches(y_positions[idx] + 0.18),
            Inches(1.75),
            Inches(0.42),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = colors[idx]
        chip.line.fill.background()
        tf = chip.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.name = "Aptos"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        box = slide.shapes.add_textbox(Inches(x_positions[idx] + 0.25), Inches(y_positions[idx] + 0.75), Inches(5.0), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = body
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT

    add_footer(slide)


def results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Results and Key Findings", "What the simulator shows")

    bullets = [
        "Small quanta improve responsiveness, but they increase context-switch overhead",
        "Larger quanta reduce switching cost and can improve turnaround on some workloads",
        "When context-switch cost increases, the best-performing quantum usually shifts upward",
        "The best quantum depends on which metric is being optimized",
        "Large quantum values cause RR behavior to approach FCFS-like execution",
    ]
    add_bullets(slide, Inches(0.75), Inches(1.45), Inches(7.3), Inches(4.7), bullets, font_size=20)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.7), Inches(4.0), Inches(4.15)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LIGHT
    tf = panel.text_frame
    lines = [
        "Example live comparisons",
        "RR sweep with `--best-by response`",
        "RR sweep with `--best-by throughput`",
        "FCFS baseline on same workload",
        "Theatre animation for intuition",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(18 if idx == 0 else 15)
        p.font.bold = idx == 0
        p.font.color.rgb = NAVY if idx == 0 else TEXT
        p.space_after = Pt(8)

    add_footer(slide)


def contribution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Why This Is More Than a Toy", "Project depth")

    cols = [
        ("Core scheduling", ["RR engine", "Context-switch overhead", "Arrival handling", "Metrics"]),
        ("Experiment platform", ["Quantum sweeps", "CSV export", "Best-by metric selection", "Synthetic workloads"]),
        ("Presentation value", ["Live animation", "Trace mode", "FCFS baseline", "Demo automation"]),
    ]
    x = [0.8, 4.55, 8.3]
    accents = [BLUE, TEAL, GOLD]
    for idx, (head, lines) in enumerate(cols):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x[idx]), Inches(1.7), Inches(3.3), Inches(4.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT

        top = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x[idx]), Inches(1.7), Inches(3.3), Inches(0.6)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = accents[idx]
        top.line.fill.background()
        tf = top.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = head
        tf.paragraphs[0].font.name = "Aptos"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        add_bullets(slide, Inches(x[idx] + 0.2), Inches(2.45), Inches(2.85), Inches(2.9), lines, font_size=17)

    add_footer(slide)


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_band(slide, "Conclusion", "Main takeaway")

    left = [
        "Round Robin quantum selection is workload-dependent",
        "No single q is universally best across all metrics",
        "Response time, turnaround, throughput, and overhead pull the decision in different directions",
        "A simulator makes those tradeoffs visible, measurable, and explainable",
    ]
    add_bullets(slide, Inches(0.8), Inches(1.55), Inches(7.2), Inches(4.9), left, font_size=21)

    quote = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.25), Inches(2.0), Inches(4.1), Inches(2.9)
    )
    quote.fill.solid()
    quote.fill.fore_color.rgb = WHITE
    quote.line.color.rgb = LIGHT
    tf = quote.text_frame
    tf.paragraphs[0].text = "Best quantum depends on:"
    tf.paragraphs[0].font.name = "Aptos"
    tf.paragraphs[0].font.size = Pt(19)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    for line in ["workload pattern", "context-switch cost", "optimization metric"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)

    add_footer(slide)


def demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, color=RGBColor(245, 248, 252))

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(5.4), Inches(1.1))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Time to Demo"
    run.font.name = "Aptos Display"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NAVY

    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.45), Inches(6.0), Inches(3.7)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = WHITE
    note.line.color.rgb = LIGHT
    tf = note.text_frame
    lines = [
        "Quick reminders",
        "Have WSL already open in the project directory",
        "Run `bash demo.sh` for the guided demo",
        "Use Ctrl+C to stop the theatre animation and continue",
        "If needed, run the animation directly with the theatre command",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(19 if idx == 0 else 15)
        p.font.bold = idx == 0
        p.font.color.rgb = NAVY if idx == 0 else TEXT
        p.space_after = Pt(8)

    cmd = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(2.25), Inches(5.15), Inches(4.2)
    )
    cmd.fill.solid()
    cmd.fill.fore_color.rgb = RGBColor(30, 35, 45)
    cmd.line.fill.background()
    tf = cmd.text_frame
    commands = [
        "WSL commands",
        "cd /mnt/c/dev/PrincOS/OSproject",
        "bash demo.sh",
        "",
        "Direct theatre mode:",
        "./rrsim --input workloads/big-mix.txt \\",
        "  --quantum 3 --cs 1 --animate \\",
        "  --style theatre --screen 110x30 --delay 40",
    ]
    for idx, line in enumerate(commands):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas" if idx > 0 else "Aptos"
        p.font.size = Pt(17 if idx == 0 else 13)
        p.font.bold = idx == 0
        p.font.color.rgb = WHITE if idx > 0 else RGBColor(224, 231, 244)
        p.space_after = Pt(4)

    add_footer(slide, text="Final slide: switch to terminal and run the live demo")


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    section_slide(
        prs,
        "Problem and Motivation",
        [
            "Round Robin is widely taught and widely used, but choosing a good time quantum is not obvious",
            "A poor quantum can improve one metric while damaging another",
            "Context-switch overhead makes the tradeoff more realistic",
            "The project goal is to turn that tradeoff into something measurable and visual",
        ],
        right_note=[
            "Project question",
            "How does quantum choice affect responsiveness, overhead, and completion behavior under different workloads?",
        ],
    )
    section_slide(
        prs,
        "Project Objectives",
        [
            "Implement a configurable Round Robin scheduler simulator in C",
            "Support arrival times, CPU burst lengths, and optional context-switch cost",
            "Compute waiting, turnaround, response, throughput, and utilization metrics",
            "Provide trace and animated terminal views of scheduling behavior",
            "Sweep across quantum values and identify strong candidates by selected metrics",
        ],
    )
    architecture_slide(prs)
    metrics_slide(prs)
    workloads_slide(prs)
    results_slide(prs)
    contribution_slide(prs)
    conclusion_slide(prs)
    demo_slide(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build_deck()
    print(OUT_PATH)
